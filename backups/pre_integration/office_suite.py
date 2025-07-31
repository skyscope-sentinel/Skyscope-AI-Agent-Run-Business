import os
import json
import logging
from typing import Dict, List, Optional, Any

# --- Dependency Checks and Imports ---
# This setup ensures the module can be imported even if optional dependencies are missing,
# providing clear guidance to the user.

try:
    import docx
    from docx.shared import Inches
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: 'python-docx' library not found. Document features will be disabled. Install with 'pip install python-docx'.")

try:
    import pptx
    from pptx.util import Inches as PptxInches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: 'python-pptx' library not found. Presentation features will be disabled. Install with 'pip install python-pptx'.")

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    print("Warning: 'openpyxl' library not found. Spreadsheet features will be disabled. Install with 'pip install openpyxl'.")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - [%(levelname)s] - %(message)s')
logger = logging.getLogger(__name__)


# --- Mock/Placeholder for AI Agent ---
class MockAgent:
    """A mock AI agent to simulate structured content generation."""
    def run(self, task: str) -> str:
        logger.info(f"MockAgent received task: {task[:100]}...")
        if "generate a document structure" in task:
            return json.dumps({
                "title": "The Future of Renewable Energy",
                "sections": [
                    {
                        "heading": "Introduction",
                        "paragraphs": [
                            "Renewable energy is at the forefront of the global conversation on climate change and sustainability.",
                            "This document explores the current landscape, key technologies, and future outlook for sources like solar, wind, and geothermal power."
                        ]
                    },
                    {
                        "heading": "Key Technologies",
                        "paragraphs": [
                            "Solar Power: Photovoltaic (PV) panels and concentrated solar power (CSP) are leading the charge.",
                            "Wind Power: Onshore and offshore wind turbines are becoming increasingly efficient and cost-effective.",
                            "Geothermal Energy: Tapping into the Earth's own heat provides a consistent and reliable power source."
                        ]
                    }
                ]
            })
        elif "generate a presentation structure" in task:
            return json.dumps({
                "slides": [
                    {"type": "title", "title": "Quarterly Business Review", "subtitle": "Q3 2025 Performance"},
                    {"type": "content", "title": "Financial Highlights", "points": ["Revenue up 15% YoY to $2.5M", "Net profit margin increased to 22%", "Reduced operational costs by 8%"]},
                    {"type": "content", "title": "Key Achievements", "points": ["Launched new 'QuantumLeap' product line", "Expanded into the European market", "Secured strategic partnership with OmniCorp"]},
                    {"type": "content", "title": "Next Quarter Goals", "points": ["Increase market share by 5%", "Hire 10 new engineers", "Begin R&D for 'Project Phoenix'"]}
                ]
            })
        elif "generate spreadsheet data" in task:
            return json.dumps({
                "headers": ["Month", "Revenue", "Expenses", "Profit"],
                "rows": [
                    ["January", 50000, 35000, 15000],
                    ["February", 55000, 37000, 18000],
                    ["March", 62000, 40000, 22000]
                ],
                "formulas": {
                    "E1": "Profit",
                    "E2": "=B2-C2" 
                }
            })
        return "{}"


class OfficeSuite:
    """
    A class to encapsulate a generative AI-powered office suite.

    This class provides methods to create and edit documents, presentations,
    and spreadsheets by interpreting natural language prompts.
    """

    def __init__(self, agent: Any = None, output_dir: str = "office_suite_output"):
        """
        Initializes the OfficeSuite.

        Args:
            agent (Any): An AI agent instance capable of processing prompts.
                         If None, a mock agent is used for demonstration.
            output_dir (str): The directory where generated files will be saved.
        """
        self.agent = agent or MockAgent()
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
        logger.info(f"OfficeSuite initialized. Output will be saved to '{self.output_dir}'.")

    # --- Document (DOCX) Methods ---

    def create_document_from_prompt(self, prompt: str, filename: str) -> Optional[str]:
        """
        Generates a Word document from a natural language prompt.

        Args:
            prompt (str): A prompt describing the desired document content.
            filename (str): The desired filename for the output .docx file.

        Returns:
            The full path to the generated document, or None on failure.
        """
        if not DOCX_AVAILABLE:
            logger.error("Cannot create document: 'python-docx' is not installed.")
            return None

        logger.info(f"Generating document content for prompt: '{prompt}'")
        ai_prompt = f"Generate a document structure for the following topic: '{prompt}'. The output must be a JSON object with a 'title' (string) and 'sections' (a list of objects, each with a 'heading' string and a 'paragraphs' list of strings)."
        
        try:
            response = self.agent.run(ai_prompt)
            content = json.loads(response)
            
            if 'title' not in content or 'sections' not in content:
                raise ValueError("AI response is missing required 'title' or 'sections' keys.")

            return self._build_docx(content, filename)
        except Exception as e:
            logger.error(f"Failed to create document: {e}")
            return None

    def _build_docx(self, content: Dict[str, Any], filename: str) -> str:
        """Builds a .docx file from structured content."""
        doc = docx.Document()
        doc.add_heading(content.get('title', 'Untitled Document'), level=0)

        for section in content.get('sections', []):
            if 'heading' in section:
                doc.add_heading(section['heading'], level=1)
            for paragraph in section.get('paragraphs', []):
                doc.add_paragraph(paragraph)
            doc.add_paragraph() # Add a space between sections

        output_path = os.path.join(self.output_dir, f"{filename}.docx")
        doc.save(output_path)
        logger.info(f"Document successfully saved to '{output_path}'")
        return output_path

    # --- Presentation (PPTX) Methods ---

    def create_presentation_from_prompt(self, prompt: str, filename: str) -> Optional[str]:
        """
        Generates a PowerPoint presentation from a natural language prompt.

        Args:
            prompt (str): A prompt describing the desired presentation content.
            filename (str): The desired filename for the output .pptx file.

        Returns:
            The full path to the generated presentation, or None on failure.
        """
        if not PPTX_AVAILABLE:
            logger.error("Cannot create presentation: 'python-pptx' is not installed.")
            return None

        logger.info(f"Generating presentation content for prompt: '{prompt}'")
        ai_prompt = f"Generate a presentation structure for the following topic: '{prompt}'. The output must be a JSON object with a 'slides' key, which is a list of objects. Each slide object should have a 'type' ('title' or 'content'), a 'title' string, and optionally a 'subtitle' string (for title slides) or a 'points' list of strings (for content slides)."
        
        try:
            response = self.agent.run(ai_prompt)
            content = json.loads(response)

            if 'slides' not in content:
                raise ValueError("AI response is missing the required 'slides' key.")

            return self._build_pptx(content, filename)
        except Exception as e:
            logger.error(f"Failed to create presentation: {e}")
            return None

    def _build_pptx(self, content: Dict[str, Any], filename: str) -> str:
        """Builds a .pptx file from structured slide data."""
        prs = pptx.Presentation()
        
        for slide_data in content.get('slides', []):
            slide_type = slide_data.get('type')
            if slide_type == 'title':
                layout = prs.slide_layouts[0] # Title slide layout
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.get('title', '')
                if 'subtitle' in slide_data:
                    slide.placeholders[1].text = slide_data.get('subtitle', '')
            elif slide_type == 'content':
                layout = prs.slide_layouts[1] # Title and Content layout
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = slide_data.get('title', '')
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.clear()
                for point in slide_data.get('points', []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
        
        output_path = os.path.join(self.output_dir, f"{filename}.pptx")
        prs.save(output_path)
        logger.info(f"Presentation successfully saved to '{output_path}'")
        return output_path

    # --- Spreadsheet (XLSX) Methods ---

    def create_spreadsheet_from_prompt(self, prompt: str, filename: str) -> Optional[str]:
        """
        Generates an Excel spreadsheet from a natural language prompt.

        Args:
            prompt (str): A prompt describing the desired spreadsheet data.
            filename (str): The desired filename for the output .xlsx file.

        Returns:
            The full path to the generated spreadsheet, or None on failure.
        """
        if not OPENPYXL_AVAILABLE:
            logger.error("Cannot create spreadsheet: 'openpyxl' is not installed.")
            return None

        logger.info(f"Generating spreadsheet data for prompt: '{prompt}'")
        ai_prompt = f"Generate spreadsheet data for the following request: '{prompt}'. The output must be a JSON object with 'headers' (a list of strings) and 'rows' (a list of lists, where each inner list represents a row)."
        
        try:
            response = self.agent.run(ai_prompt)
            content = json.loads(response)

            if 'headers' not in content or 'rows' not in content:
                raise ValueError("AI response is missing required 'headers' or 'rows' keys.")

            return self._build_xlsx(content, filename)
        except Exception as e:
            logger.error(f"Failed to create spreadsheet: {e}")
            return None

    def _build_xlsx(self, content: Dict[str, Any], filename: str) -> str:
        """Builds an .xlsx file from structured table data."""
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"

        headers = content.get('headers', [])
        ws.append(headers)

        for row_data in content.get('rows', []):
            ws.append(row_data)
        
        # Example of handling simple formulas if provided by AI
        if 'formulas' in content:
            for cell, formula in content['formulas'].items():
                try:
                    ws[cell] = formula
                except Exception as e:
                    logger.warning(f"Could not apply formula to cell {cell}: {e}")

        output_path = os.path.join(self.output_dir, f"{filename}.xlsx")
        wb.save(output_path)
        logger.info(f"Spreadsheet successfully saved to '{output_path}'")
        return output_path


if __name__ == '__main__':
    logger.info("--- OfficeSuite Demonstration ---")
    
    # Initialize the suite with the mock agent
    suite = OfficeSuite()
    
    # 1. Create a document
    logger.info("\n--- Creating a Word Document ---")
    doc_prompt = "A brief report on the future of renewable energy, covering solar and wind power."
    doc_path = suite.create_document_from_prompt(doc_prompt, "Renewable_Energy_Report")
    if doc_path:
        logger.info(f"✅ Document created at: {doc_path}")
    else:
        logger.error("❌ Document creation failed.")
        
    # 2. Create a presentation
    logger.info("\n--- Creating a PowerPoint Presentation ---")
    ppt_prompt = "A quarterly business review presentation with 4 slides: title, financial highlights, key achievements, and next quarter goals."
    ppt_path = suite.create_presentation_from_prompt(ppt_prompt, "Q3_Business_Review")
    if ppt_path:
        logger.info(f"✅ Presentation created at: {ppt_path}")
    else:
        logger.error("❌ Presentation creation failed.")
        
    # 3. Create a spreadsheet
    logger.info("\n--- Creating an Excel Spreadsheet ---")
    xls_prompt = "A simple financial summary spreadsheet with columns for Month, Revenue, Expenses, and Profit for the first quarter."
    xls_path = suite.create_spreadsheet_from_prompt(xls_prompt, "Q1_Financials")
    if xls_path:
        logger.info(f"✅ Spreadsheet created at: {xls_path}")
    else:
        logger.error("❌ Spreadsheet creation failed.")

```
